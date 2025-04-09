import sys
import os
import time
import random
import logging
from PyQt6.QtWidgets import (
    QApplication, QWidget, QVBoxLayout, QLabel, QPushButton,
    QFileDialog, QGridLayout, QMessageBox, QCheckBox, QInputDialog,QGraphicsDropShadowEffect,
    QProgressDialog, QWidget, QLabel, QPushButton,QLineEdit,QFileDialog, QMessageBox,
)
from PyQt6.QtCore import Qt
from PyQt6.QtGui import QIcon, QPixmap, QCursor,QColor,QPixmap, QImage
from PyQt6.QtCore import Qt, QTimer, QSettings
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
import tempfile
import pdfplumber
import pandas as pd
import pikepdf
from docx2pdf import convert
import subprocess
import shutil



class PDFToolApp(QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("PDF Tool App")
        self.setGeometry(100, 100, 1100, 700)

        self.settings = QSettings("PhonePe", "PDF_App")
        self.is_dark_mode = self.settings.value("dark_mode", True, type=bool)  # Dark mode default

        self.light_stylesheet = self.build_stylesheet(dark=False)
        self.dark_stylesheet = self.build_stylesheet(dark=True)
        self.setStyleSheet(self.dark_stylesheet if self.is_dark_mode else self.light_stylesheet)

        self.layout = QVBoxLayout()
        self.setLayout(self.layout)

        # Top Label
        self.top_label = QLabel("PHONEPE Loves PDF")
        self.top_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.top_label.setStyleSheet("font-size: 24px; font-weight: bold; margin-bottom: 15px;")
        self.layout.addWidget(self.top_label)

        self.theme_toggle = QCheckBox("Dark Mode")
        self.theme_toggle.setChecked(self.is_dark_mode)
        self.theme_toggle.stateChanged.connect(self.toggle_theme)
        self.layout.addWidget(self.theme_toggle, alignment=Qt.AlignmentFlag.AlignRight)

        self.grid = QGridLayout()
        self.grid.setSpacing(20)
        self.layout.addLayout(self.grid)

        self.animations = {}

        tools = [
            ("📎 Merge PDF", "icons/merge.png", self.merge_pdfs, "Combine multiple PDF files into one."),
            ("✂️ Split PDF", "icons/split.png", self.split_pdf, "Split a PDF into individual pages."),
            ("📉 Compress PDF", "icons/compress.png", self.compress_pdf, "Reduce PDF file size."),
            ("📝 PDF to Word", "icons/pdf_to_word.png", self.pdf_to_word, "Convert PDF to Word document."),
            ("📊 PDF to PowerPoint", "icons/pdf_to_ppt.png", self.pdf_to_ppt, "Convert PDF to PowerPoint slides."),
            ("📈 PDF to Excel", "icons/pdf_to_excel.png", self.pdf_to_excel, "Convert PDF to Excel sheet."),
            ("🖼️ PDF to JPG", "icons/pdf_to_jpg.png", self.pdf_to_jpg, "Convert PDF pages to images."),
            ("🗑️ Remove PDF Pages ", "icons/pdf_to_jpg.png", self.remove_pdf_pages, "Convert PDF pages to images."),
            ("💧 Watermark", "icons/watermark.png", self.add_watermark_to_pdf, "Add watermark to PDF."),
            ("🖊️ Unlcok PDF", "icons/sign.png", self.unlock_pdf, "Digitally sign a PDF document."),
            ("🖊️ Lcok PDF", "icons/sign.png", self.lock_pdf_with_password, "Digitally sign a PDF document."),
            ("🔍 OCR PDF", "icons/ocr.png", self.ocr_pdf, "Extract text from scanned PDF using OCR."),
            ("🛠️ Repair PDF", "icons/repair.png", self.repair_pdf, "Fix corrupted or damaged PDFs."),
             ("✏️ Edit PDF", "icons/edit.png", self.edit_pdf_in_word, "Edit text or images in PDF."),
        ]


        positions = [(i, j) for i in range(5) for j in range(4)]
        for position, (label, icon_path, func, tooltip) in zip(positions, tools):
            btn = QPushButton(label)
            btn.setCursor(QCursor(Qt.CursorShape.PointingHandCursor))
            btn.setToolTip(tooltip)
            if os.path.exists(icon_path):
                btn.setIcon(QIcon(icon_path))
                btn.setIconSize(QPixmap(icon_path).rect().size())
            btn.clicked.connect(lambda checked, b=btn, f=func: self.animate_click(b, f))
            btn.setFixedSize(220, 110)
            btn.installEventFilter(self)
            self.grid.addWidget(btn, *position)

        self.footer = QLabel("Developed by Chandan S")
        self.footer.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.footer.setStyleSheet("font-size: 14px; margin-top: 20px;")
        self.layout.addWidget(self.footer)

    def toggle_theme(self):
        self.is_dark_mode = self.theme_toggle.isChecked()
        self.setStyleSheet(self.dark_stylesheet if self.is_dark_mode else self.light_stylesheet)
        self.settings.setValue("dark_mode", self.is_dark_mode)

    def build_stylesheet(self, dark=False):
        if dark:
            return """
                QWidget { background-color: #1e1e2e; color: #f8f8f2; font-size: 16px; }
                QPushButton {
                    background-color: #2c2c3c;
                    color: #f8f8f2;
                    border: 2px solid #444;
                    border-radius: 15px;
                    padding: 10px;
                }
                QPushButton:hover {
                    background-color: #3c3c5c;
                    border: 2px solid #6272a4;
                    color: #ff79c6;
                }
            """
        else:
            return """
                QWidget { background-color: #f4f6f9; color: #2c2c2c; font-size: 16px; }
                QPushButton {
                    background-color: #ffffff;
                    color: #2c2c2c;
                    border: 2px solid #d1d1d1;
                    border-radius: 15px;
                    padding: 10px;
                }
                QPushButton:hover {
                    background-color: #e6f0ff;
                    border: 2px solid #3399ff;
                    color: #0055cc;
                }
            """

    def animate_click(self, button, func):
        original_text = button.text()
        button.setText("Working...")
        QTimer.singleShot(150, lambda: self.run_action(button, func, original_text))

    def run_action(self, button, func, original_text):
        func()
        button.setText(original_text)

    def eventFilter(self, obj, event):
        if isinstance(obj, QPushButton):
            if event.type() == event.Type.Enter:
                self.animate_button(obj, enter=True)
            elif event.type() == event.Type.Leave:
                self.animate_button(obj, enter=False)
        return super().eventFilter(obj, event)

    def animate_button(self, button, enter=True):
        if enter:
            shadow = QGraphicsDropShadowEffect(self)
            shadow.setBlurRadius(15)
            shadow.setXOffset(0)
            shadow.setYOffset(0)
            shadow.setColor(QColor("#3399ff"))  # Glow color
            button.setGraphicsEffect(shadow)
        else:
            button.setGraphicsEffect(None)

    def merge_pdfs(self):
        files, _ = QFileDialog.getOpenFileNames(self, "Select PDFs to Merge", "", "PDF Files (*.pdf)")
        if not files:
            return  # User cancelled

        try:
            merger = PdfMerger()
            progress_dialog = QProgressDialog("Merging PDFs...", "Cancel", 0, len(files), self)
            progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)
            progress_dialog.setWindowTitle("Merging PDFs")

            for i, file in enumerate(files):
                if progress_dialog.wasCanceled():
                    merger.close()
                    return  # User cancelled
                # Simple PDF validation (can be improved)
                if not file.lower().endswith(".pdf") or not os.path.isfile(file):
                    QMessageBox.warning(self, "Warning", f"Invalid file: {file}")
                    continue
                merger.append(file)
                progress_dialog.setValue(i + 1)

            output_file, _ = QFileDialog.getSaveFileName(self, "Save Merged PDF", "merged.pdf", "PDF Files (*.pdf)")
            if output_file:
                merger.write(output_file)
                merger.close()
                QMessageBox.information(self, "Success", "PDFs merged successfully!")

        except Exception as e:
            QMessageBox.critical(self, "Error", f"An error occurred during PDF merging: {e}")

        finally:
            if 'merger' in locals():  # Ensure merger is defined
                merger.close()

      
    def split_pdf(self):
        logging.debug("split_pdf function started.")  # Log start of function

        file, _ = QFileDialog.getOpenFileName(self, "Select PDF to Split", "", "PDF Files (*.pdf)")
        if not file:
            logging.debug("User cancelled file selection.")
            return

        logging.debug(f"Selected file: {file}")

        # Step 1: Ask for page ranges input
        try:
            reader = PdfReader(file)
            total_pages = len(reader.pages)
            logging.debug(f"Total pages in PDF: {total_pages}")
        except Exception as e:
            logging.error(f"Error reading PDF: {e}", exc_info=True)
            QMessageBox.critical(self, "Error", f"Error reading PDF: {e}")
            return

        input_text, ok = QInputDialog.getText(self, "Enter Pages", f"Enter page ranges to split (1 to {total_pages}, e.g. 1-3,5,7):")
        if not ok or not input_text.strip():
            logging.debug("User cancelled or entered invalid page ranges.")
            QMessageBox.warning(self, "No Input", "No page ranges entered.")
            return

        def parse_ranges(ranges_text):
            pages = set()
            for part in ranges_text.split(','):
                part = part.strip()
                if '-' in part:
                    try:
                        start, end = part.split('-')
                        pages.update(range(int(start) - 1, int(end)))
                    except ValueError:
                        logging.warning(f"Invalid range format: {part}")
                else:
                    try:
                        pages.add(int(part) - 1)
                    except ValueError:
                        logging.warning(f"Invalid page number: {part}")
            return sorted(p for p in pages if 0 <= p < total_pages)

        selected_pages = parse_ranges(input_text)
        if not selected_pages:
            logging.debug("No valid page numbers selected.")
            QMessageBox.warning(self, "Invalid Pages", "No valid page numbers found.")
            return

        logging.debug(f"Selected pages: {selected_pages}")

        # Step 2: Ask for output folder
        output_dir = QFileDialog.getExistingDirectory(self, "Select Output Directory")
        if not output_dir:
            logging.debug("User cancelled directory selection.")
            return

        logging.debug(f"Output directory: {output_dir}")

        # Step 3: Split and save
        base_name = os.path.splitext(os.path.basename(file))[0]
        logging.debug(f"Base filename: {base_name}")

        progress_dialog = QProgressDialog("Splitting PDF...", "Cancel", 0, len(selected_pages), self)
        progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)
        progress_dialog.setWindowTitle("Splitting PDF")

        try:
            for i, page_num in enumerate(selected_pages):
                if progress_dialog.wasCanceled():
                    logging.debug(f"User cancelled splitting at page {page_num}.")
                    return

                writer = PdfWriter()
                writer.add_page(reader.pages[page_num])

                # Generate a truly unique filename
                timestamp = int(time.time())
                random_id = random.randint(1000, 9999)
                output_filename = f"{base_name}_page_{page_num + 1}_{timestamp}_{random_id}.pdf"
                output_path = os.path.join(output_dir, output_filename)

                logging.debug(f"Saving page {page_num + 1} to: {output_path}")

                try:
                    with open(output_path, "wb") as f:
                        writer.write(f)  # **CRITICAL: This is the mistake!**
                    writer.close()  # Explicitly close the writer
                    logging.debug(f"Page {page_num + 1} saved successfully.")
                except Exception as write_error:
                    logging.error(f"Error writing page {page_num + 1}: {write_error}", exc_info=True)

                progress_dialog.setValue(i + 1)

            QMessageBox.information(self, "Done", f"Split {len(selected_pages)} pages successfully.")
            logging.info("PDF split successfully.")

        except Exception as main_error:
            logging.error(f"An unexpected error occurred: {main_error}", exc_info=True)
            QMessageBox.critical(self, "Error", f"An error occurred during PDF splitting: {main_error}")


    def compress_pdf(self):
        file, _ = QFileDialog.getOpenFileName(self, "Select PDF to Compress", "", "PDF Files (*.pdf)")
        if not file:
            return

        output_file, _ = QFileDialog.getSaveFileName(self, "Save Compressed PDF", "compressed.pdf", "PDF Files (*.pdf)")
        if not output_file:
            return

        try:
            doc = fitz.open(file)

            # Just save with compress=True (let MuPDF do its job)
            doc.save(output_file, garbage=4, deflate=True)
            doc.close()

            QMessageBox.information(self, "Success", "PDF compressed successfully!")
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Compression failed: {e}")


    def pdf_to_word(self):
        input_file, _ = QFileDialog.getOpenFileName(self, "Select PDF File", "", "PDF Files (*.pdf)")
        if input_file:
            output_file, _ = QFileDialog.getSaveFileName(self, "Save Word File", "converted.docx", "Word Files (*.docx)")
            if output_file:
                try:
                    cv = Converter(input_file)
                    cv.convert(output_file, start=0, end=None)
                    cv.close()
                    QMessageBox.information(self, "Success", "PDF successfully converted to Word document.")
                except Exception as e:
                    QMessageBox.warning(self, "Error", f"Failed to convert PDF to Word.\n{str(e)}")

    def pdf_to_ppt(self):
        # Step 1: Ask for PDF
        pdf_path, _ = QFileDialog.getOpenFileName(self, "Select PDF to Convert", "", "PDF Files (*.pdf)")
        if not pdf_path:
            return

        # Step 2: Ask for Output Path
        ppt_path, _ = QFileDialog.getSaveFileName(self, "Save PowerPoint File", "converted.pptx", "PowerPoint Files (*.pptx)")
        if not ppt_path:
            return

        # Step 3: Convert PDF pages to images
        doc = fitz.open(pdf_path)
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]  # no title/content

        for page in doc:
            pix = page.get_pixmap(dpi=150)
            temp_img_path = os.path.join(tempfile.gettempdir(), f"page_{page.number}.png")
            pix.save(temp_img_path)

            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(temp_img_path, Inches(0), Inches(0), width=prs.slide_width)

        prs.save(ppt_path)
        QMessageBox.information(self, "Success", f"PPT saved to:\n{ppt_path}")
    
       
    def pdf_to_excel(self):
        # Step 1: Ask user to select a PDF file
        pdf_path, _ = QFileDialog.getOpenFileName(self, "Select PDF to Convert", "", "PDF Files (*.pdf)")
        if not pdf_path:
            return

        # Step 2: Ask for output Excel file path
        excel_path, _ = QFileDialog.getSaveFileName(self, "Save Excel File", "converted.xlsx", "Excel Files (*.xlsx)")
        if not excel_path:
            return

        table_dfs = []
        lines_data = []

        try:
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    # Extract tables
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            df = pd.DataFrame(table[1:], columns=table[0])
                            df["Page"] = page_num
                            table_dfs.append(df)

                    # Extract lines of text
                    lines = page.extract_text().split('\n') if page.extract_text() else []
                    for line_num, line in enumerate(lines, 1):
                        lines_data.append({
                            "Page": page_num,
                            "Line Number": line_num,
                            "Text": line
                        })

            # Save both tables and lines to Excel
            with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
                if table_dfs:
                    pd.concat(table_dfs, ignore_index=True).to_excel(writer, sheet_name="Extracted Tables", index=False)
                else:
                    pd.DataFrame([{"Note": "No tables found."}]).to_excel(writer, sheet_name="Extracted Tables", index=False)

                pd.DataFrame(lines_data).to_excel(writer, sheet_name="All Text Lines", index=False)

            QMessageBox.information(self, "Success", f"Excel saved to:\n{excel_path}")

        except Exception as e:
            QMessageBox.critical(self, "Error", f"An error occurred:\n{e}")
    
    def pdf_to_jpg(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Select PDF File", "", "PDF Files (*.pdf)")
        if not file_path:
            return

        output_dir = QFileDialog.getExistingDirectory(self, "Select Output Folder")
        if not output_dir:
            return

        quality, ok = QInputDialog.getInt(self, "Image Quality", "Enter image quality (1–100):", 80, 1, 100)
        if not ok:
            return

        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))

                # Convert pixmap to PNG bytes, then to QImage
                img_data = pix.tobytes("png")
                qimg = QImage.fromData(img_data)

                output_path = os.path.join(output_dir, f"page_{page_num + 1}.jpg")
                qimg.save(output_path, "JPG", quality)

            QMessageBox.information(self, "Success", f"PDF pages converted to JPGs in:\n{output_dir}")
        except Exception as e:
            logging.error(f"Error converting PDF to JPG: {e}", exc_info=True)
            QMessageBox.critical(self, "Error", f"Failed to convert PDF: {e}")

    def ocr_pdf(self):
        file, _ = QFileDialog.getOpenFileName(self, "Select PDF for OCR", "", "PDF Files (*.pdf)")
        if file:
            output_txt, _ = QFileDialog.getSaveFileName(self, "Save Extracted Text", "ocr_text.txt", "Text Files (*.txt)")
            if output_txt:
                doc = fitz.open(file)
                text = ""
                for page in doc:
                    img = page.get_pixmap()
                    img_pil = Image.frombytes("RGB", [img.width, img.height], img.samples)
                    text += pytesseract.image_to_string(img_pil) + "\n"
                with open(output_txt, "w", encoding="utf-8") as f:
                    f.write(text)

    def repair_pdf(self):
        file, _ = QFileDialog.getOpenFileName(self, "Select PDF to Repair", "", "PDF Files (*.pdf)")
        if file:
            output_file, _ = QFileDialog.getSaveFileName(self, "Save Repaired PDF", "repaired.pdf", "PDF Files (*.pdf)")
            if output_file:
                reader = PdfReader(file)
                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)
                with open(output_file, "wb") as out_file:
                    writer.write(out_file)

    def remove_pdf_pages(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Open PDF to Remove Pages", "", "PDF Files (*.pdf)")
        if not file_path:
            return

        doc = fitz.open(file_path)

        total_pages = len(doc)
        pages_str, ok = QInputDialog.getText(
            self, "Remove Pages", f"Enter page numbers to remove (1-{total_pages}, comma-separated):"
        )

        if not ok or not pages_str.strip():
            return

        try:
            pages_to_remove = sorted(set(int(p.strip()) - 1 for p in pages_str.split(",") if p.strip().isdigit()), reverse=True)
            for page_num in pages_to_remove:
                if 0 <= page_num < len(doc):
                    doc.delete_page(page_num)

            output_path, _ = QFileDialog.getSaveFileName(self, "Save Modified PDF", "modified.pdf", "PDF Files (*.pdf)")
            if output_path:
                doc.save(output_path)
                QMessageBox.information(self, "Success", f"Pages removed successfully!\nSaved to:\n{output_path}")
        except Exception as e:
            logging.error(f"Error removing pages: {e}", exc_info=True)
            QMessageBox.critical(self, "Error", f"Failed to remove pages: {e}")

    def add_watermark_to_pdf(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Select PDF to Add Watermark", "", "PDF Files (*.pdf)")
        if not file_path:
            return

        text, ok = QInputDialog.getText(self, "Watermark Text", "Enter watermark text:")
        if not ok or not text:
            return

        try:
            doc = fitz.open(file_path)
            for page in doc:
                rect = page.rect
                page.insert_textbox(
                    rect,
                    text,
                    fontsize=60,
                    color=(0.8, 0.8, 0.8),  # light gray
                    rotate=90,  # fixed to a valid value
                    align=fitz.TEXT_ALIGN_CENTER,
                    overlay=True
                )

            save_path, _ = QFileDialog.getSaveFileName(self, "Save Watermarked PDF", "watermarked.pdf", "PDF Files (*.pdf)")
            if save_path:
                doc.save(save_path)
                QMessageBox.information(self, "Success", f"Watermark added and saved to:\n{save_path}")

        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to add watermark:\n{e}")

    def unlock_pdf(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Select PDF to Unlock", "", "PDF Files (*.pdf)")
        if not file_path:
            return

        password, ok = QInputDialog.getText(self, "Enter Password", "PDF is encrypted. Enter password:", QLineEdit.EchoMode.Password)
        if not ok or not password:
            return

        try:
            doc = fitz.open(file_path)
            if not doc.is_encrypted:
                QMessageBox.information(self, "Info", "PDF is not encrypted.")
                return

            if not doc.authenticate(password):
                QMessageBox.critical(self, "Failed", "Incorrect password.")
                return

            save_path, _ = QFileDialog.getSaveFileName(self, "Save Unlocked PDF", "unlocked.pdf", "PDF Files (*.pdf)")
            if save_path:
                doc.save(save_path, encryption=fitz.PDF_ENCRYPT_NONE)
                QMessageBox.information(self, "Success", f"Unlocked PDF saved to:\n{save_path}")

        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to unlock PDF:\n{e}")

    def lock_pdf_with_password(self):
        input_path, _ = QFileDialog.getOpenFileName(self, "Select PDF to Lock", "", "PDF Files (*.pdf)")
        if not input_path:
            return

        password, ok = QInputDialog.getText(self, "Set Password", "Enter password:", QLineEdit.EchoMode.Password)
        if not ok or not password:
            return

        output_path, _ = QFileDialog.getSaveFileName(self, "Save Locked PDF", "locked.pdf", "PDF Files (*.pdf)")
        if not output_path:
            return

        try:
            with pikepdf.open(input_path) as pdf:
                pdf.save(
                    output_path,
                    encryption=pikepdf.Encryption(
                        owner=password,
                        user=password,
                        R=6  # AES-256 encryption
                    )
                )
            QMessageBox.information(self, "Success", "PDF locked successfully!")
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to lock PDF:\n{e}")

    def edit_pdf_in_word(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Open PDF", "", "PDF Files (*.pdf)")
        if not file_path:
            return

        try:
            # 1. Convert PDF to DOCX
            temp_dir = tempfile.mkdtemp()
            docx_path = os.path.join(temp_dir, "converted.docx")
            pdf_path = os.path.abspath(file_path)

            cv = Converter(pdf_path)
            cv.convert(docx_path, start=0, end=None)
            cv.close()

            # 2. Open in Word for editing
            subprocess.run(["start", "", docx_path], shell=True)

            QMessageBox.information(self, "Edit in Word", "Make changes and save the Word file. Click OK when done.")

            # 3. Convert edited DOCX back to PDF
            output_pdf_path, _ = QFileDialog.getSaveFileName(self, "Save Edited PDF", "edited.pdf", "PDF Files (*.pdf)")
            if output_pdf_path:
                convert(docx_path, output_pdf_path)
                QMessageBox.information(self, "Success", f"Edited PDF saved to:\n{output_pdf_path}")

        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to edit PDF in Word:\n{e}")
        finally:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
                
    def placeholder(self):
        QMessageBox.information(self, "Coming Soon", "This feature is not implemented yet.")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = PDFToolApp()
    window.show()
    sys.exit(app.exec())